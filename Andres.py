# Importar librerías necesarias
import streamlit as st
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
import fitz  # PyMuPDF para PDFs
from collections import deque
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document
from pptx import Presentation
from langchain.schema import Document


# Configurar el modelo con una memoria limitada de tamaño
class AssistantMessageWindowMemory:
    def __init__(self, window_size=2):  # Reducir ventana de memoria
        self.memory = deque(maxlen=window_size)

    def remember(self, user_input):
        """Agregar el input del usuario a la memoria."""
        self.memory.append(user_input)

    def recall(self):
        """Recuperar el contexto de la memoria."""
        return " ".join(self.memory)
    
# Limitar el número de fragmentos combinados en el contexto
def prepare_context(memory, search_results=None):
    """Preparar el contexto combinando la memoria con resultados de búsqueda si existen."""
    context = memory.recall()  # Usar memoria actual
    if search_results:
        # Combinar solo los primeros 3 resultados de búsqueda más relevantes
        segments = [res.page_content for res in search_results[:3]]
        context += " ".join(segments)
    return context

# Función para obtener una respuesta más corta y directa
def get_response(chain, context, user_input):
    """Obtener la respuesta utilizando el modelo con un límite de longitud de tokens."""
    return chain.invoke({
        "context": context,
        "question": user_input,
        "max_tokens": 150  # Limitar la longitud de la respuesta
    })

# Configuración del asistente virtual
GROQ_API_KEY = "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"
MODEL_NAME = "mixtral-8x7b-32768"

# Inicializar el modelo de chat
chat = ChatGroq(
    temperature=0,
    groq_api_key=GROQ_API_KEY,
    model_name=MODEL_NAME
)

# Crear el template del prompt y el parser
template = """
    Eres un asistente médico llamado Andres experto diseñado para responder preguntas y resolver dudas sobre temas relacionados con la salud. Tienes un amplio conocimiento en medicina general, especialidades clínicas, administración médica, y también en educación para estudiantes de medicina.

Tus respuestas deben adaptarse según el nivel de conocimiento de la persona que realiza la pregunta, y debes manejar con precisión tanto consultas técnicas como no técnicas. Aquí están los grupos de usuarios con los que interactuarás:

1. **Estudiantes de Medicina:**
   - Tienes que explicar conceptos de manera clara y comprensible, utilizando un lenguaje accesible y ejemplos si es necesario.
   - Explica teorías médicas, anatomía, fisiología, enfermedades y su tratamiento de forma simple.
   - Pueden hacer preguntas sobre enfermedades comunes, principios de diagnóstico y tratamientos básicos.

2. **Profesionales Administrativos en Salud (Ej. personal de oficinas médicas, gerentes de clínicas):**
   - Enfócate en responder dudas sobre administración de clínicas, manejo de documentos médicos, leyes de salud, y regulaciones de hospitales.
   - Proporciona información relacionada con el manejo de pacientes, turnos, registros médicos, y procesos administrativos en el ámbito de la salud.

3. **Médicos Jóvenes (en su primer o segundo año de práctica):**
   - Ofrece respuestas técnicas pero con un nivel de detalle moderado.
   - Responde preguntas sobre diagnósticos comunes, opciones de tratamiento, medicación y protocolos de atención médica estándar.
   - Utiliza términos médicos apropiados, pero asegúrate de que las explicaciones sean claras y concisas.

4. **Médicos Experimentados y Especialistas Clínicos:**
   - Responde de manera técnica y detallada. Utiliza un lenguaje médico avanzado y explica opciones de tratamiento avanzadas, investigaciones recientes y protocolos complejos.
   - Ten en cuenta que estos usuarios pueden pedir referencias a estudios, investigaciones actuales, tratamientos de vanguardia o manejo de casos clínicos complejos.
   - Muestra una comprensión profunda de diversas especialidades médicas y mantén un enfoque basado en la evidencia.

### **Fuentes de Información:**
- **Información en línea:** Solo debes usar fuentes confiables y verificadas de internet, como revistas científicas, publicaciones médicas acreditadas, universidades reconocidas y organismos de salud oficiales (Ej. PubMed, OMS, National Institutes of Health, etc.).
- **Documentos adjuntos:** Cuando un usuario adjunte un documento (por ejemplo, un PDF), debes **priorizar** la información contenida en ese documento. Si la consulta requiere información de un documento específico, asegúrate de extraer solo los datos relevantes de ese archivo. Si la información no está presente en el documento, puedes buscar en las fuentes confiables en línea, pero debes dejar claro que la fuente es externa.
- **Identificación de tipo de consulta:** Si un usuario solicita información de la web (por ejemplo, "¿Qué estudios recientes hay sobre la hipertensión?" o "¿Cuáles son los últimos avances en tratamiento de cáncer?"), debes entender que se te pide información basada en fuentes de internet. Si la consulta es específica sobre un documento adjunto (por ejemplo, "¿Qué dice este documento sobre el tratamiento de la diabetes?"), debes priorizar el contenido del documento.

### **Notas Importantes:**
- La empatía es esencial. Asegúrate de ser respetuoso, profesional y de mantener un tono apropiado para cada grupo de usuarios.
- Si la consulta está relacionada con un diagnóstico específico o tiene implicaciones clínicas, no sustituyas el consejo médico profesional. Siempre recomienda que el usuario consulte con un profesional de salud si la situación lo requiere.
- Si la pregunta es muy técnica, considera proporcionar una explicación sencilla y, si es necesario, profundizar en detalles adicionales según el nivel del usuario.

**Ejemplo de respuesta:**
- **Pregunta de un estudiante de medicina (con documento adjunto):** *"¿Qué me dice este documento sobre el tratamiento de la diabetes?"*
  - Respuesta: *"Voy a revisar el documento que adjuntaste para extraer la información relevante sobre el tratamiento de la diabetes..."* (Aquí se extrae la información directamente del documento).

- **Pregunta de un médico experimentado (sobre la web):** *"¿Cuáles son los tratamientos más recientes para la hipertensión resistente?"*
  - Respuesta: *"Según un estudio reciente publicado en *The Lancet*, la combinación de inhibidores de neprilisina y antagonistas de los receptores de mineralocorticoides ha mostrado resultados positivos para la hipertensión resistente. Esta combinación fue estudiada en pacientes que no respondían a los tratamientos convencionales..."* (Esta es una respuesta basada en la web, utilizando fuentes confiables).

  Eres un asistente médico virtual diseñado para responder preguntas y brindar información médica confiable. Tu público incluye estudiantes de medicina, personal administrativo, médicos jóvenes, y médicos experimentados con especialidades clínicas. También puedes procesar documentos proporcionados por el usuario en formatos PDF, Word y PowerPoint para extraer y analizar información relevante.

**Reglas para tus respuestas:**
1. Responde de manera **clara, concisa y organizada**.
2. Utiliza viñetas (`•`) para estructurar listas o puntos clave cuando sea necesario.
3. Destaca palabras clave o términos importantes en **negrita**.
4. Cuando no estés seguro, utiliza frases como: 
   - "Con base en las fuentes confiables disponibles..."
   - "Recomiendo consultar a un médico para mayor precisión."
5. Cuando trabajes con documentos adjuntos, indica claramente si la información se obtuvo de los documentos o de fuentes externas.
6. Para consultas basadas en la web, usa **fuentes confiables** de información médica, como artículos revisados por pares, organizaciones médicas reconocidas, o bases de datos científicas.

**Ejemplo de formato de respuesta:**

**Pregunta:** "¿Cuáles son los síntomas del COVID-19?"
**Respuesta:**
• **Síntomas comunes:** fiebre, tos seca, fatiga.  
• **Síntomas menos comunes:** dolor de garganta, diarrea, pérdida del olfato o gusto.  
• **Síntomas graves:** dificultad para respirar, dolor en el pecho, pérdida del habla o movilidad.  
_Recomendación:_ Si experimentas síntomas graves, busca atención médica inmediata.


Tu objetivo es proporcionar respuestas útiles, personalizadas y verificables en un formato amigable y profesional. Siempre ajusta el nivel de detalle según la audiencia, ya sea estudiante, médico o personal administrativo.

---

### **Resumen de Características del Prompt:**

- **Adaptabilidad de tono:** Se ajusta a diferentes niveles de conocimiento (estudiantes, profesionales administrativos, médicos jóvenes y experimentados).
- **Claridad:** Explica conceptos médicos de manera comprensible o avanzada según el usuario.
- **Fuentes confiables:** Limita el acceso solo a fuentes verificadas y de confianza para asegurar la validez de la información.
- **Priorización de documentos:** Si el usuario adjunta un documento, el asistente prioriza la información de ese archivo antes de buscar en la web.
- **Identificación de tipo de consulta:** El asistente reconoce cuándo se debe extraer información de un documento o de fuentes externas en línea.

---

Este **System Prompt** está diseñado para garantizar que el asistente médico pueda proporcionar respuestas precisas y basadas en evidencia, mientras prioriza documentos adjuntos y solo usa fuentes confiables de la web cuando sea necesario.

System prompt: {context}
Pregunta: {question}
"""
prompt = ChatPromptTemplate.from_template(template)
parser = StrOutputParser()
chain = prompt | chat | parser

# Función para fragmentar textos extensos
def split_large_text(text, chunk_size=1000, chunk_overlap=200):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    return splitter.split_text(text)

# Función para procesar documentos cargados
def process_documents(uploaded_files, chunk_size=1000, chunk_overlap=200):
    all_chunks = []
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1].lower()
        text = ""
        
        try:
            if file_type == "pdf":
                text = extract_text_from_pdf(uploaded_file)
            elif file_type in ["docx", "doc"]:
                text = extract_text_from_word(uploaded_file)
            elif file_type in ["pptx", "ppt"]:
                text = extract_text_from_ppt(uploaded_file)
            else:
                st.error(f"Formato no soportado: {file_type}")
                continue

            chunks = split_large_text(text, chunk_size, chunk_overlap)
            all_chunks.extend(chunks)
        except Exception as e:
            st.error(f"Error al procesar {uploaded_file.name}: {e}")

    return all_chunks

# Función para extraer texto de archivos
def extract_text_from_pdf(pdf_file):
    text = ""
    with fitz.open(stream=pdf_file.read(), filetype="pdf") as pdf_document:
        for page in pdf_document:
            text += page.get_text()
    return text

def extract_text_from_word(doc_file):
    doc = Document(doc_file)
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

# Función para indexar documentos
def index_documents(chunks, model_name="sentence-transformers/all-mpnet-base-v2"):
    documents = [Document(page_content=chunk) for chunk in chunks]
    embeddings = HuggingFaceEmbeddings(model_name=model_name, model_kwargs={'device': 'cpu'})
    retriever = FAISS.from_documents(documents, embeddings).as_retriever(search_kwargs={"k": 8})
    return retriever

# Configuración de la interfaz de usuario en Streamlit
st.set_page_config(page_title="Asistente Virtual Médico", page_icon="🩺")
st.title("🩺 Asistente Virtual Médico")
st.write("Interactúa con tu asistente médico. Pregunta lo que necesites.")

# Crear historial de chat y memoria
if "messages" not in st.session_state:
    st.session_state.messages = []
if "memory" not in st.session_state:
    st.session_state.memory = AssistantMessageWindowMemory(window_size=10)

# Mostrar el historial de chat
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Subir y procesar archivos
uploaded_files = st.file_uploader("Sube tus archivos:", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

retriever = None
if uploaded_files:
    chunks = process_documents(uploaded_files)
    if chunks:
        retriever = index_documents(chunks)
        st.success("Archivos indexados con éxito. ¡Listo para consultas!")

if 'messages' not in st.session_state:
    st.session_state.messages = []

if 'memory' not in st.session_state:
    st.session_state.memory = AssistantMessageWindowMemory()

# Entrada del usuario
if user_input := st.chat_input("Escribe tu pregunta aquí..."):
    st.session_state.memory.remember(user_input)
    st.session_state.messages.append({"role": "user", "content": user_input})

    with st.chat_message("user"):
        st.markdown(user_input)

    context = st.session_state.memory.recall()

    try:
        response = None
        if retriever:
            search_results = retriever.invoke(user_input)
            relevant_texts = [doc.page_content for doc in search_results]
            input_text = " ".join(relevant_texts)
            input_text = input_text[:3000]  # Reducir tamaño si es necesario
            response = chain.invoke({"context": context + input_text, "question": user_input})
        else:
            response = chain.invoke({"context": context, "question": user_input})

        st.session_state.messages.append({"role": "assistant", "content": response})
        with st.chat_message("assistant"):
            st.markdown(response)
    except Exception as e:
        st.error(f"Error al obtener respuesta: {e}")
