# Importar librerías necesarias
import streamlit as st
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
import fitz  # PyMuPDF para PDFs
from collections import deque
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document  # Para leer archivos Word
from pptx import Presentation  # Para leer archivos PowerPoint
from langchain.schema import Document
from dotenv import load_dotenv
import os

# Cargar las variables del archivo .env
load_dotenv()

# Obtener las variables
groq_api_key = os.getenv('GROQ_API_KEY')
model_name = os.getenv('MODEL_NAME')

# Verificar si las variables están bien cargadas
if not groq_api_key or not model_name:
    print("Error: API key o nombre de modelo no están configurados correctamente.")
    exit()

# Inicializar el modelo de chat
try:
    chat = ChatGroq(
        temperature=0,  # Configuración de la temperatura para la generación de respuestas
        groq_api_key=groq_api_key,
        model_name=model_name
    )
    print("Modelo de chat inicializado correctamente.")
except Exception as e:
    print(f"Error al inicializar el modelo: {e}")

# Función para preparar el contexto con menos datos
def prepare_context(memory, search_results=None):
    """Prepara el contexto combinando solo una parte de la memoria y los resultados de búsqueda relevantes"""
    context = memory.recall()  # Recupera la memoria
    if search_results:
        # Limitar a los primeros 2 resultados más relevantes
        segments = [res.page_content for res in search_results[:2]]  # Usamos solo 2 resultados
        context += " ".join(segments)
    return context


# Función para indexar los documentos
def index_documents(pages, model_name="sentence-transformers/all-mpnet-base-v2"):
    # Crear una lista de objetos Document
    documents = [Document(page_content=page) for page in pages]  # Usa "page_content" en lugar de "content"
    
    # Crear embeddings con HuggingFace
    hf = HuggingFaceEmbeddings(model_name=model_name, model_kwargs={'device': 'cpu'})
    
    # Indexar documentos con FAISS
    faiss = FAISS.from_documents(documents, hf)
    return faiss.as_retriever(search_kwargs={"k": 5})

# Función para extraer texto de un archivo PDF
def extract_text_from_pdf(pdf_file):
    text = ""
    with fitz.open(stream=pdf_file.read(), filetype="pdf") as pdf_document:
        for page in pdf_document:
            text += page.get_text()
    return text

# Función para extraer texto de un archivo Word (.docx)
def extract_text_from_word(doc_file):
    doc = Document(doc_file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Función para extraer texto de una presentación PowerPoint (.pptx)
def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Función para procesar los documentos cargados
def process_documents(uploaded_files, chunk_size=2000, chunk_overlap=200):
    all_pages = []
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,  # Ajusta el tamaño del fragmento
        chunk_overlap=chunk_overlap,  # Ajusta el solapamiento
        length_function=len,  # Usar la longitud de caracteres para dividir
        is_separator_regex=False,  # No usar un separador regex
    )
    
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1].lower()
        text = ""
        
        # Procesar cada tipo de archivo
        if file_type == "pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif file_type in ["docx", "doc"]:
            text = extract_text_from_word(uploaded_file)
        elif file_type in ["pptx", "ppt"]:
            text = extract_text_from_ppt(uploaded_file)
        else:
            st.error(f"Formato de archivo no soportado: {file_type}")
            continue
        
        if text:
            # Dividir el texto en fragmentos
            chunks = text_splitter.split_text(text)

            # Agregar los fragmentos a la lista
            all_pages.extend(chunks)
    
    return all_pages


# Configuración de la interfaz de usuario en Streamlit
st.set_page_config(page_title="Asistente Virtual Médico", page_icon="🩺")
st.title("🩺 Asistente Virtual Médico")
st.write("Interactúa con tu asistente médico personal. Pregunta lo que necesites.")

# Crear o recuperar el historial de chat y la memoria
if "messages" not in st.session_state:
    st.session_state.messages = []
if "memory" not in st.session_state:
    st.session_state.memory = AssistantMessageWindowMemory(window_size=10)

# Mostrar el historial de chat
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# **Cambio aquí:** Utilizar un buscador de archivos con file uploader para PDF, Word y PPT
uploaded_files = st.file_uploader("Selecciona los archivos PDF, Word o PowerPoint con historiales clínicos:", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

# Procesar los archivos cargados y realizar la indexación RAG
retriever = None
if uploaded_files:
    all_pages = process_documents(uploaded_files)
    
    if all_pages:
        retriever = index_documents(all_pages)
        st.success("Archivos procesados y documentos indexados con éxito. Ahora puedes hacer preguntas sobre su contenido.")
    else:
        st.warning("No se encontraron archivos válidos para procesar.")

# Obtener la entrada del usuario y el contexto
if user_input := st.chat_input("Escribe tu pregunta aquí..."):
    # Agregar la pregunta a la memoria
    st.session_state.memory.remember(user_input)
    # Agregar la pregunta al historial de mensajes
    st.session_state.messages.append({"role": "user", "content": user_input})

    # Mostrar la pregunta del usuario
    with st.chat_message("user"):
        st.markdown(user_input)

    # Obtener los resultados de búsqueda
    search_results = retriever.invoke(user_input) if retriever else None
    
    # Aquí definimos el contexto con memoria + resultados de búsqueda
    context = prepare_context(st.session_state.memory, search_results)

    # Asegúrate de no exceder el límite de tokens
    max_tokens = 5000

    try:
        # Realizar la llamada al modelo usando el contexto
        response = None
        if retriever:
            # Si estás usando el recuperador (retriever), añades los fragmentos obtenidos
            search_results = retriever.invoke(user_input)
            segments = [i.page_content for i in search_results]
            response = chain.invoke({"context": context + " ".join(segments), "question": user_input, "max_tokens": max_tokens})
        else:
            response = chain.invoke({"context": context, "question": user_input, "max_tokens": max_tokens})

        # Agregar la respuesta al historial de mensajes
        st.session_state.messages.append({"role": "assistant", "content": response})

        # Mostrar la respuesta
        with st.chat_message("assistant"):
            st.markdown(response)

    except Exception as e:
        st.error(f"Error al obtener respuesta: {e}")
